"""
生成湘潭大学风格答辩PPTX - 简约版
"""

import os
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
FIGURES_DIR = os.path.join(BASE_DIR, "outputs", "figures")
OUTPUT_DIR = os.path.join(BASE_DIR, "outputs", "presentation")

# 简约配色
MAROON = RGBColor(0x7A, 0x0E, 0x1E)
DARK = RGBColor(0x2D, 0x2D, 0x2D)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT = RGBColor(0xF5, 0xF5, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BORDER = RGBColor(0xE0, 0xE0, 0xE0)
ACCENT = RGBColor(0x1A, 0x5C, 0x8A)

W = Emu(12192000)
H = Emu(6858000)


def tb(slide, l, t, w, h, text, sz=14, color=DARK, bold=False,
       align=PP_ALIGN.LEFT, font="微软雅黑", ls=1.3):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    p.line_spacing = Pt(int(sz * ls))
    return box


def ml(slide, l, t, w, h, lines, sz=12, color=DARK, bold=False,
       font="微软雅黑", ls=1.4, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(sz)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font
        p.alignment = align
        p.line_spacing = Pt(int(sz * ls))
    return box


def bullets(slide, l, t, w, h, items, sz=12, color=DARK):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"  {item}"
        p.font.size = Pt(sz)
        p.font.color.rgb = color
        p.font.name = "微软雅黑"
        p.line_spacing = Pt(int(sz * 1.7))
    return box


def img(slide, name, l, t, w=None):
    path = os.path.join(FIGURES_DIR, name)
    if os.path.exists(path):
        kw = {"image_file": path, "left": l, "top": t}
        if w:
            kw["width"] = w
        slide.shapes.add_picture(**kw)
        return True
    return False


def rect(slide, l, t, w, h, fill, line_color=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line_color:
        s.line.color.rgb = line_color
        s.line.width = Pt(0.5)
    else:
        s.line.fill.background()
    return s


def num_tag(slide, l, t, num):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, Emu(320000), Emu(320000))
    s.fill.solid()
    s.fill.fore_color.rgb = MAROON
    s.line.fill.background()
    tb(slide, l, t + Emu(50000), Emu(320000), Emu(220000), num,
       sz=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)


def metric(slide, l, t, w, val, label, vc=MAROON):
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, Emu(480000))
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT
    box.line.color.rgb = BORDER
    box.line.width = Pt(0.5)
    tb(slide, l + Emu(30000), t + Emu(40000), w - Emu(60000), Emu(240000),
       val, sz=20, color=vc, bold=True, align=PP_ALIGN.CENTER)
    tb(slide, l + Emu(30000), t + Emu(290000), w - Emu(60000), Emu(150000),
       label, sz=9, color=GRAY, align=PP_ALIGN.CENTER)


def bottom_line(slide):
    rect(slide, Emu(450000), Emu(6400000), Emu(11300000), Emu(15000), MAROON)


def page(slide, num):
    tb(slide, Emu(11400000), Emu(6450000), Emu(500000), Emu(200000),
       str(num), sz=9, color=GRAY, align=PP_ALIGN.RIGHT)


# ============ 内容页模板 ============

def make_slide(prs, title, num, build_fn):
    slide = prs.slides.add_slide(prs.slide_layouts[0])

    # 左侧窄色条
    rect(slide, Emu(0), Emu(0), Emu(60000), H, MAROON)

    # 标题
    tb(slide, Emu(450000), Emu(300000), Emu(8000000), Emu(400000),
       title, sz=22, color=MAROON, bold=True)

    # 标题下细线
    rect(slide, Emu(450000), Emu(720000), Emu(2000000), Emu(12000), MAROON)

    build_fn(slide)

    bottom_line(slide)
    page(slide, num)
    return slide


# ============ 封面 ============

def cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])

    # 顶部色块
    rect(slide, Emu(0), Emu(0), W, Emu(120000), MAROON)

    # 左侧竖线装饰
    rect(slide, Emu(800000), Emu(1400000), Emu(60000), Emu(2600000), MAROON)

    # 校名
    tb(slide, Emu(1100000), Emu(1600000), Emu(4000000), Emu(500000),
       "湘潭大学", sz=32, color=MAROON, bold=True)

    # 课程
    tb(slide, Emu(1100000), Emu(2200000), Emu(4000000), Emu(300000),
       "大数据分析与数据挖掘 · 课程设计答辩", sz=14, color=GRAY)

    # 主标题
    tb(slide, Emu(1100000), Emu(3000000), Emu(10000000), Emu(600000),
       "基于机器学习与强化学习的\n股票交易策略建模及量化风险评估研究",
       sz=24, color=DARK, bold=True)

    # 分隔线
    rect(slide, Emu(1100000), Emu(3900000), Emu(3000000), Emu(10000), BORDER)

    # 信息
    info = [
        "学    院：自动化与电子信息学院",
        "专    业：人工智能",
        "姓    名：【待补充】",
        "学    号：【待补充】",
        "指导教师：【待补充】",
    ]
    ml(slide, Emu(1100000), Emu(4100000), Emu(5000000), Emu(1800000),
       info, sz=13, color=DARK, ls=1.8)

    # 日期
    tb(slide, Emu(9000000), Emu(5800000), Emu(2500000), Emu(300000),
       "2026 年 6 月", sz=12, color=GRAY, align=PP_ALIGN.RIGHT)

    # 底部色块
    rect(slide, Emu(0), Emu(6738000), W, Emu(120000), MAROON)

    page(slide, 1)
    return slide


# ============ 目录 ============

def toc(prs):
    def build(slide):
        items = [
            "研究背景与任务目标",
            "总体技术路线",
            "数据获取与预处理",
            "特征工程",
            "机器学习模型结果",
            "LSTM 模型结果",
            "回测与风险评估",
            "强化学习与 Q-learning",
            "总结与展望",
        ]
        for i, text in enumerate(items):
            y = Emu(1100000) + i * Emu(560000)
            num_tag(slide, Emu(600000), y, f"{i+1:02d}")
            tb(slide, Emu(1100000), y + Emu(40000), Emu(8000000), Emu(280000),
               text, sz=15, color=DARK, bold=True)
            if i < len(items) - 1:
                rect(slide, Emu(1100000), y + Emu(400000),
                     Emu(9000000), Emu(8000), BORDER)

    return make_slide(prs, "目  录", 2, build)


# ============ 研究背景 ============

def s_background(prs):
    def build(slide):
        tb(slide, Emu(450000), Emu(950000), Emu(5000000), Emu(350000),
           "为什么研究这个问题？", sz=16, color=MAROON, bold=True)
        bullets(slide, Emu(450000), Emu(1400000), Emu(5000000), Emu(1500000), [
            "金融时间序列具有高噪声、非线性和时序性",
            "短期方向信号较弱，复杂模型不一定优于简单基线",
            "收益表现必须结合最大回撤与尾部风险共同解释",
        ], sz=12)

        tb(slide, Emu(450000), Emu(3100000), Emu(5000000), Emu(350000),
           "本项目要完成什么？", sz=16, color=MAROON, bold=True)
        bullets(slide, Emu(450000), Emu(3550000), Emu(5000000), Emu(1500000), [
            "数据准备：获取、清洗并检查 SPY 日线数据",
            "模型预测：比较逻辑回归、随机森林与 LSTM",
            "风险评估：计算回撤、Sharpe、VaR 与 CVaR",
            "强化学习：完成简化环境与 Q-learning 流程",
        ], sz=12)

        # 右侧标注
        box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Emu(7000000), Emu(1200000),
            Emu(4500000), Emu(3800000))
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT
        box.line.color.rgb = BORDER
        box.line.width = Pt(0.5)

        tb(slide, Emu(7300000), Emu(1500000), Emu(3900000), Emu(300000),
           "实验定位", sz=14, color=MAROON, bold=True, align=PP_ALIGN.CENTER)
        tb(slide, Emu(7300000), Emu(2100000), Emu(3900000), Emu(2500000),
           "用于课程方法验证与风险意识训练\n不构成任何实际投资建议\n\n"
           "以 SPY ETF 为研究对象\n覆盖 2015-2025 年日线数据\n"
           "完成 ML / DL / RL 三层方法对比",
           sz=12, color=DARK, align=PP_ALIGN.CENTER)

    return make_slide(prs, "研究背景与任务目标", 3, build)


# ============ 技术路线 ============

def s_pipeline(prs):
    def build(slide):
        steps = [
            ("01", "数据获取", "Yahoo Finance"),
            ("02", "数据清洗", "质量检查"),
            ("03", "特征工程", "9 项技术特征"),
            ("04", "ML 预测", "LR / RF"),
            ("05", "LSTM", "20 日窗口"),
            ("06", "风险评估", "Sharpe / VaR"),
            ("07", "Q-learning", "状态/动作/奖励"),
            ("08", "综合分析", "结论与边界"),
        ]

        bw, bh = Emu(1300000), Emu(600000)
        gap = Emu(160000)
        sx = Emu(450000)

        for i, (num, title, desc) in enumerate(steps):
            col, row = i % 4, i // 4
            x = sx + col * (bw + gap)
            y = Emu(1000000) + row * (bh + Emu(500000))

            box = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, x, y, bw, bh)
            box.fill.solid()
            box.fill.fore_color.rgb = LIGHT
            box.line.color.rgb = MAROON if row == 0 else ACCENT
            box.line.width = Pt(1)

            tb(slide, x + Emu(20000), y + Emu(50000),
               bw - Emu(40000), Emu(180000), num,
               sz=11, color=MAROON if row == 0 else ACCENT, bold=True,
               align=PP_ALIGN.CENTER)
            tb(slide, x + Emu(20000), y + Emu(220000),
               bw - Emu(40000), Emu(180000), title,
               sz=12, color=DARK, bold=True, align=PP_ALIGN.CENTER)
            tb(slide, x + Emu(20000), y + Emu(400000),
               bw - Emu(40000), Emu(150000), desc,
               sz=9, color=GRAY, align=PP_ALIGN.CENTER)

            if i < len(steps) - 1 and col < 3:
                arr = slide.shapes.add_shape(
                    MSO_SHAPE.RIGHT_ARROW,
                    x + bw + Emu(15000), y + Emu(230000),
                    Emu(130000), Emu(130000))
                arr.fill.solid()
                arr.fill.fore_color.rgb = BORDER
                arr.line.fill.background()

        # 底部原则
        tb(slide, Emu(450000), Emu(4200000), Emu(10500000), Emu(600000),
           "关键原则：按时间顺序划分训练/测试 | 信号用于下一期收益 | 如实呈现所有结果",
           sz=11, color=DARK, bold=True)
        tb(slide, Emu(450000), Emu(4700000), Emu(10500000), Emu(400000),
           "Python 3.10  |  yfinance  |  pandas  |  scikit-learn  |  PyTorch  |  matplotlib",
           sz=10, color=GRAY)

    return make_slide(prs, "总体技术路线", 4, build)


# ============ 数据获取 ============

def s_data(prs):
    def build(slide):
        metric(slide, Emu(450000), Emu(950000), Emu(2800000), "2766", "交易日")
        metric(slide, Emu(3500000), Emu(950000), Emu(2800000), "6", "核心字段")
        metric(slide, Emu(6550000), Emu(950000), Emu(2800000), "2015-2025", "覆盖区间")

        tb(slide, Emu(450000), Emu(1700000), Emu(4000000), Emu(300000),
           "预处理步骤", sz=15, color=MAROON, bold=True)
        bullets(slide, Emu(450000), Emu(2100000), Emu(4500000), Emu(1800000), [
            "字段筛选：Open / High / Low / Close / Adj Close / Volume",
            "缺失值处理：前向填充 + 后向填充",
            "日期去重、异常检查、时间排序",
        ], sz=11)

        img(slide, "SPY_close_price.png", Emu(5800000), Emu(1700000), w=Emu(5500000))

        tb(slide, Emu(450000), Emu(5600000), Emu(10000000), Emu(400000),
           "数据质量：日期单调递增，无重复；核心字段缺失值为 0",
           sz=10, color=GRAY, bold=True)

    return make_slide(prs, "数据获取与预处理", 5, build)


# ============ 特征工程 ============

def s_features(prs):
    def build(slide):
        tb(slide, Emu(450000), Emu(950000), Emu(10000000), Emu(300000),
           "在 OHLCV 基础上构建 9 项技术特征", sz=13, color=GRAY)

        cats = [
            ("收益变化", "return_1d / return_5d", "短期与短周期累计变化"),
            ("趋势信息", "ma_5 / 10 / 20 / 60", "平滑噪声，描述价格趋势"),
            ("风险与位置", "volatility_20 / volume_ma_20\nclose_ma20_ratio",
             "波动、活跃度与相对位置"),
        ]

        cw = Emu(3000000)
        for i, (cat, feats, desc) in enumerate(cats):
            x = Emu(450000) + i * (cw + Emu(250000))
            y = Emu(1400000)

            box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, cw, Emu(800000))
            box.fill.solid()
            box.fill.fore_color.rgb = LIGHT
            box.line.color.rgb = MAROON
            box.line.width = Pt(0.8)

            tb(slide, x + Emu(100000), y + Emu(50000), cw - Emu(200000), Emu(180000),
               cat, sz=12, color=MAROON, bold=True, align=PP_ALIGN.CENTER)
            tb(slide, x + Emu(100000), y + Emu(250000), cw - Emu(200000), Emu(200000),
               feats, sz=10, color=DARK, bold=True, align=PP_ALIGN.CENTER)
            tb(slide, x + Emu(100000), y + Emu(500000), cw - Emu(200000), Emu(200000),
               desc, sz=9, color=GRAY, align=PP_ALIGN.CENTER)

        img(slide, "SPY_feature_correlation_heatmap.png",
            Emu(450000), Emu(2600000), w=Emu(5500000))

        tb(slide, Emu(6500000), Emu(2800000), Emu(4800000), Emu(2000000),
           "滚动窗口产生初始缺失值\n删除后：2707 行 x 15 列\n\n"
           "OHLC 和均线高度相关\n因为它们都来自价格水平",
           sz=12, color=DARK)

    return make_slide(prs, "特征工程", 6, build)


# ============ ML 结果 ============

def s_ml(prs):
    def build(slide):
        metric(slide, Emu(450000), Emu(950000), Emu(3500000),
               "0.5886", "逻辑回归 Accuracy", vc=ACCENT)
        metric(slide, Emu(4500000), Emu(950000), Emu(3500000),
               "0.4188", "随机森林 Accuracy", vc=RGBColor(0x99, 0x33, 0x00))

        tb(slide, Emu(450000), Emu(1700000), Emu(4500000), Emu(300000),
           "逻辑回归", sz=14, color=MAROON, bold=True)
        bullets(slide, Emu(450000), Emu(2100000), Emu(4500000), Emu(1000000), [
            "recall = 1.0，542 个样本全部预测为上涨",
            "准确率 = 上涨类别占比 (58.86%)",
            "存在明显类别偏向",
        ], sz=11)

        tb(slide, Emu(450000), Emu(3300000), Emu(4500000), Emu(300000),
           "随机森林", sz=14, color=MAROON, bold=True)
        bullets(slide, Emu(450000), Emu(3700000), Emu(4500000), Emu(1000000), [
            "上涨 recall = 0.0282，仅预测少量上涨日",
            "当前特征和参数下泛化不足",
            "accuracy 必须结合 precision / recall / F1 阅读",
        ], sz=11)

        img(slide, "SPY_model_prediction_comparison.png",
            Emu(5800000), Emu(1700000), w=Emu(5500000))

        tb(slide, Emu(450000), Emu(5500000), Emu(10500000), Emu(400000),
           "结论：逻辑回归相对较好，但不能只看 Accuracy",
           sz=11, color=ACCENT, bold=True)

    return make_slide(prs, "机器学习模型结果", 7, build)


# ============ LSTM ============

def s_lstm(prs):
    def build(slide):
        metric(slide, Emu(450000), Emu(950000), Emu(3500000),
               "0.4540", "LSTM Accuracy", vc=ACCENT)
        metric(slide, Emu(4500000), Emu(950000), Emu(3500000),
               "0.4916", "LSTM AUC", vc=RGBColor(0x99, 0x33, 0x00))

        tb(slide, Emu(450000), Emu(1700000), Emu(4500000), Emu(300000),
           "实验设置", sz=14, color=MAROON, bold=True)
        bullets(slide, Emu(450000), Emu(2100000), Emu(4500000), Emu(1500000), [
            "输入：连续 20 个交易日的技术特征窗口",
            "输出：下一交易日涨跌标签 (二分类)",
            "模型：单层 LSTM (hidden=32)",
            "训练：30 轮，BCELoss + Adam (lr=0.001)",
            "测试集：522 个序列样本",
        ], sz=11)

        tb(slide, Emu(450000), Emu(3900000), Emu(4500000), Emu(300000),
           "结果解释", sz=14, color=MAROON, bold=True)
        bullets(slide, Emu(450000), Emu(4300000), Emu(4500000), Emu(1200000), [
            "AUC 接近 0.5，区分能力接近随机",
            "深度模型不一定优于简单基线",
            "金融短期预测噪声较大",
        ], sz=11)

        img(slide, "SPY_lstm_training_curve.png",
            Emu(5800000), Emu(1700000), w=Emu(5500000))

    return make_slide(prs, "LSTM 模型结果", 8, build)


# ============ 回测与风险 ============

def s_backtest(prs):
    def build(slide):
        metric(slide, Emu(450000), Emu(950000), Emu(2600000),
               "12.54%", "SPY 年化收益", vc=RGBColor(0x00, 0x66, 0x00))
        metric(slide, Emu(3300000), Emu(950000), Emu(2600000),
               "-34.10%", "SPY 最大回撤", vc=RGBColor(0xCC, 0x00, 0x00))
        metric(slide, Emu(6150000), Emu(950000), Emu(2600000),
               "0.7019", "Sharpe Ratio", vc=ACCENT)

        tb(slide, Emu(450000), Emu(1700000), Emu(4500000), Emu(300000),
           "风险指标", sz=14, color=MAROON, bold=True)
        bullets(slide, Emu(450000), Emu(2100000), Emu(4500000), Emu(1500000), [
            "年化波动率：17.86%",
            "Calmar Ratio：0.3676",
            "VaR 95%：-1.68%",
            "CVaR 95%：-2.73%",
        ], sz=11)

        tb(slide, Emu(450000), Emu(3700000), Emu(4500000), Emu(300000),
           "简化 RF 策略回测", sz=14, color=MAROON, bold=True)
        bullets(slide, Emu(450000), Emu(4100000), Emu(4500000), Emu(1200000), [
            "策略：RF 预测上涨时持有，否则空仓",
            "信号 shift(1) 避免未来函数",
            "回测 Sharpe：1.1981 | 最大回撤：-0.87%",
            "未扣交易成本、滑点或税费",
        ], sz=11)

        img(slide, "SPY_strategy_equity_curve.png",
            Emu(5800000), Emu(1700000), w=Emu(5500000))

    return make_slide(prs, "回测与风险评估", 9, build)


# ============ RL ============

def s_rl(prs):
    def build(slide):
        tb(slide, Emu(450000), Emu(950000), Emu(4500000), Emu(300000),
           "交易环境设计", sz=14, color=MAROON, bold=True)
        bullets(slide, Emu(450000), Emu(1350000), Emu(4500000), Emu(1500000), [
            "状态：return_1d / close_ma20_ratio / volatility_20 离散化",
            "动作：0=空仓，1=持有 SPY",
            "奖励：下一期持仓收益 - 仓位变化成本 (0.1%)",
            "状态空间：3x3x3 = 27 个离散状态",
        ], sz=11)

        tb(slide, Emu(450000), Emu(2900000), Emu(10500000), Emu(300000),
           "三策略对比", sz=14, color=MAROON, bold=True)

        strats = [
            ("Random", "-57.64%", RGBColor(0xCC, 0x00, 0x00)),
            ("Buy & Hold", "+227.12%", RGBColor(0x00, 0x66, 0x00)),
            ("Q-learning", "+6.56%", ACCENT),
        ]
        for i, (name, ret, color) in enumerate(strats):
            x = Emu(450000) + i * Emu(3500000)
            y = Emu(3300000)
            box = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, x, y, Emu(3100000), Emu(700000))
            box.fill.solid()
            box.fill.fore_color.rgb = LIGHT
            box.line.color.rgb = color
            box.line.width = Pt(1)

            tb(slide, x, y + Emu(50000), Emu(3100000), Emu(200000),
               name, sz=12, color=color, bold=True, align=PP_ALIGN.CENTER)
            tb(slide, x, y + Emu(280000), Emu(3100000), Emu(300000),
               ret, sz=22, color=color, bold=True, align=PP_ALIGN.CENTER)

        img(slide, "SPY_rl_qlearning_equity.png",
            Emu(450000), Emu(4300000), w=Emu(5500000))

        tb(slide, Emu(6500000), Emu(4400000), Emu(4800000), Emu(1500000),
           "Q-learning 未超过买入持有\n\n"
           "原因：研究期 SPY 长期上涨\n"
           "仅 3 个离散化特征\n"
           "动作只有空仓和持有\n"
           "训练仅 50 轮\n\n"
           "价值在于完成 RL 流程并如实呈现局限",
           sz=11, color=DARK)

    return make_slide(prs, "强化学习与 Q-learning", 10, build)


# ============ 总结 ============

def s_conclusion(prs):
    def build(slide):
        tb(slide, Emu(450000), Emu(950000), Emu(10500000), Emu(300000),
           "核心结论", sz=16, color=MAROON, bold=True)
        bullets(slide, Emu(450000), Emu(1400000), Emu(10500000), Emu(2000000), [
            "完成从数据获取到强化学习实验的完整流程",
            "逻辑回归相对较好，但存在明显类别偏向",
            "LSTM AUC 接近 0.5，短期区分能力有限",
            "Q-learning 获得正收益 (+6.56%)，但显著弱于买入持有 (+227.12%)",
            "风险指标显示：长期上涨仍伴随显著下行风险 (最大回撤 -34.10%)",
        ], sz=12)

        tb(slide, Emu(450000), Emu(3500000), Emu(10500000), Emu(300000),
           "可拓展方向", sz=16, color=MAROON, bold=True)

        ext = [
            ("01", "更多特征", "宏观 / 情绪 / 基本面"),
            ("02", "多资产", "组合风险分散"),
            ("03", "滚动训练", "严格样本外验证"),
            ("04", "真实约束", "成本 / 滑点 / 执行"),
            ("05", "高级 RL", "DQN / PPO"),
        ]
        for i, (num, title, desc) in enumerate(ext):
            x = Emu(450000) + i * Emu(2100000)
            y = Emu(3950000)
            box = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, x, y, Emu(1900000), Emu(800000))
            box.fill.solid()
            box.fill.fore_color.rgb = LIGHT
            box.line.color.rgb = BORDER
            box.line.width = Pt(0.5)

            tb(slide, x, y + Emu(40000), Emu(1900000), Emu(180000),
               num, sz=14, color=MAROON, bold=True, align=PP_ALIGN.CENTER)
            tb(slide, x, y + Emu(250000), Emu(1900000), Emu(200000),
               title, sz=11, color=DARK, bold=True, align=PP_ALIGN.CENTER)
            tb(slide, x, y + Emu(480000), Emu(1900000), Emu(250000),
               desc, sz=9, color=GRAY, align=PP_ALIGN.CENTER)

        tb(slide, Emu(450000), Emu(5600000), Emu(10500000), Emu(350000),
           "本项目仅用于课程实验，不构成投资建议",
           sz=11, color=MAROON, bold=True)

    return make_slide(prs, "总结与展望", 11, build)


# ============ 致谢 ============

def thanks(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])

    rect(slide, Emu(0), Emu(0), W, Emu(120000), MAROON)

    tb(slide, Emu(0), Emu(2400000), W, Emu(700000),
       "感谢聆听", sz=36, color=MAROON, bold=True, align=PP_ALIGN.CENTER)

    tb(slide, Emu(0), Emu(3300000), W, Emu(400000),
       "敬请指导", sz=18, color=GRAY, align=PP_ALIGN.CENTER)

    rect(slide, Emu(4800000), Emu(4000000), Emu(2600000), Emu(10000), MAROON)

    tb(slide, Emu(0), Emu(4300000), W, Emu(300000),
       "基于机器学习与强化学习的股票交易策略建模及量化风险评估研究",
       sz=11, color=GRAY, align=PP_ALIGN.CENTER)

    rect(slide, Emu(0), Emu(6738000), W, Emu(120000), MAROON)

    page(slide, 12)
    return slide


# ============ 主函数 ============

def main():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    cover(prs)
    toc(prs)
    s_background(prs)
    s_pipeline(prs)
    s_data(prs)
    s_features(prs)
    s_ml(prs)
    s_lstm(prs)
    s_backtest(prs)
    s_rl(prs)
    s_conclusion(prs)
    thanks(prs)

    out = os.path.join(OUTPUT_DIR, "湘潭大学-大数据课程设计答辩-股票量化RL.pptx")
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    prs.save(out)
    print(f"Saved: {out}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
